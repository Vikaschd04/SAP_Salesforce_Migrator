#!/usr/bin/env python
"""Build the H2A Migrator executive deck as a native, editable A4-landscape PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── palette ──
BLACK=RGBColor(0x0E,0x0E,0x0E); INK=RGBColor(0x22,0x24,0x27); GRAY=RGBColor(0x53,0x56,0x5A)
FAINT=RGBColor(0x8f,0x92,0x96); LINE=RGBColor(0xd2,0xd2,0xcf); PANEL=RGBColor(0xf3,0xf3,0xf1)
WHITE=RGBColor(0xFF,0xFF,0xFF); GREEN=RGBColor(0x86,0xBC,0x25); GDEEP=RGBColor(0x04,0x6a,0x38)
TEAL=RGBColor(0x00,0x7c,0xb0); COPPER=RGBColor(0xa9,0x62,0x2f); RISK=RGBColor(0xb3,0x28,0x2d)
WARN=RGBColor(0x9a,0x6b,0x0a)
SANS="Arial"; MONO="Courier New"

# ── geometry (A4 landscape) ──
PW,PH = 11.69, 8.27
MX = 0.66
CW = PW-2*MX
EY_Y=0.52; TITLE_Y=0.9; RULE_Y=1.82; BODY=2.08
STRIP_Y=7.06; STRIP_H=0.78; FOOT_Y=7.98

prs=Presentation(); prs.slide_width=Inches(PW); prs.slide_height=Inches(PH)

def slide(bg=WHITE):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.shadow.inherit=False; r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background()
    return s

def rect(s,l,t,w,h,fill=None,ln=None,lw=0.75,shape=MSO_SHAPE.RECTANGLE,round=0.0):
    sp=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h)); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if ln is None: sp.line.fill.background()
    else: sp.line.color.rgb=ln; sp.line.width=Pt(lw)
    if shape==MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0]=round
        except Exception: pass
    return sp

def tf_box(s,l,t,w,h,anchor=MSO_ANCHOR.TOP,wrap=True):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return tf

def _track(run,pts): run._r.get_or_add_rPr().set('spc',str(int(pts*100)))

def para(tf, runs, align=PP_ALIGN.LEFT, size=11, color=INK, bold=False, name=SANS,
         italic=False, sa=3, sb=0, line=1.08):
    p0=tf.paragraphs[0]
    p = p0 if (not p0.runs) else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(sa); p.space_before=Pt(sb); p.line_spacing=line
    if isinstance(runs,str): runs=[(runs,{})]
    for txt,st in runs:
        r=p.add_run(); r.text=txt; f=r.font
        f.size=Pt(st.get('size',size)); f.bold=st.get('bold',bold); f.italic=st.get('italic',italic)
        f.name=st.get('name',name); f.color.rgb=st.get('color',color)
        if st.get('track'): _track(r,st['track'])
    return p

def eyebrow(s,kick,dark=False):
    rect(s,MX,EY_Y+0.02,0.13,0.13,fill=GREEN,shape=MSO_SHAPE.OVAL)
    tf=tf_box(s,MX+0.27,EY_Y-0.04,CW-0.3,0.32)
    para(tf,[(kick.upper(),dict(size=10.5,bold=True,color=(FAINT if dark else GRAY),name=MONO,track=2.0))])

def title(s,lines,size=29,dark=False,y=TITLE_Y):
    ink=WHITE if dark else BLACK
    tf=tf_box(s,MX,y,CW,1.05)
    for i,ln in enumerate(lines):
        runs=[(ln,dict(size=size,bold=True,color=ink,name=SANS))]
        if i==len(lines)-1: runs.append((".",dict(size=size,bold=True,color=GREEN,name=SANS)))
        para(tf,runs,size=size,line=1.02,sa=0)

def footer(s,page,dark=False):
    col=RGBColor(0x6a,0x6d,0x70) if dark else FAINT
    tf=tf_box(s,MX,FOOT_Y,4,0.25); para(tf,[("H2A MIGRATOR",dict(size=8,color=col,name=MONO,track=1.2))])
    tf2=tf_box(s,PW-MX-2.2,FOOT_Y,2.0,0.25); para(tf2,[(f"{page:02d} / 20",dict(size=8,color=col,name=MONO,track=1.2))],align=PP_ALIGN.RIGHT)
    rect(s,PW-MX-0.14,FOOT_Y+0.025,0.09,0.09,fill=GREEN,shape=MSO_SHAPE.OVAL)

def header(s,kick,lines,page,rule=True,tsize=29):
    eyebrow(s,kick); title(s,lines,size=tsize)
    if rule: rect(s,MX,RULE_Y,CW,0.02,fill=LINE)
    footer(s,page)

def card(s,l,t,w,h,top=BLACK,fill=WHITE,border=LINE):
    r=rect(s,l,t,w,h,fill=fill,ln=border,lw=0.75)
    if top: rect(s,l,t,w,0.055,fill=top)
    return r

def strip(s,label,body_runs):
    rect(s,MX,STRIP_Y,CW,STRIP_H,fill=PANEL)
    rect(s,MX,STRIP_Y,0.07,STRIP_H,fill=GREEN)
    tf=tf_box(s,MX+0.28,STRIP_Y+0.12,CW-0.55,STRIP_H-0.22)
    para(tf,[(label.upper(),dict(size=8.5,bold=True,color=GDEEP,name=MONO,track=1.4))],sa=4)
    para(tf,body_runs,size=11.5,color=INK,line=1.16)

def chips(s,items,y,l=MX):
    x=l
    for it in items:
        w=0.36+len(it['t'])*0.086
        rect(s,x,y,w,0.38,fill=WHITE,ln=LINE,lw=0.75)
        tf=tf_box(s,x+0.15,y,w-0.3,0.38,anchor=MSO_ANCHOR.MIDDLE,wrap=False)
        para(tf,it['runs'],size=9,name=MONO,color=GRAY,line=1.0)
        x+=w+0.2

def cols(n,gap,l=MX,w=CW):
    cw=(w-(n-1)*gap)/n
    return [l+i*(cw+gap) for i in range(n)], cw

# ══════════════════════════════════════════════════════════════════════════
# 1 · COVER
s=slide()
rect(s,7.9,0.7,4.9,4.9,ln=GREEN,lw=1.5,shape=MSO_SHAPE.OVAL)     # ring
rect(s,9.2,3.4,3.2,3.2,ln=LINE,lw=1.0,shape=MSO_SHAPE.OVAL)
rect(s,7.62,1.02,0.32,0.32,fill=GREEN,shape=MSO_SHAPE.OVAL)
eyebrow(s,"Executive Reference Deck · v0.7.0")
tf=tf_box(s,MX,1.4,7.6,1.6)
para(tf,[("SAP Hybris to Salesforce.",dict(size=34,bold=True,color=BLACK))],line=1.05,sa=0)
para(tf,[("Migrated by AI. Proven to run",dict(size=34,bold=True,color=BLACK)),(".",dict(size=34,bold=True,color=GREEN))],line=1.05,sa=0)
rect(s,MX,2.85,0.95,0.08,fill=GREEN)
tf=tf_box(s,MX,3.15,7.1,1.1)
para(tf,"A complete walkthrough — the problem, the architecture, how the AI agent team works, how every output is verified, and what it means for the business. One deck, every answer.",size=13.5,color=GRAY,line=1.3)
# flow line
fx=MX
for t,c,arrow in [("SAP HYBRIS · JAVA",COPPER,True),("AI AGENT TEAM",BLACK,True),("SALESFORCE · APEX",TEAL,False)]:
    w=0.3+len(t)*0.083
    rect(s,fx,4.55,w,0.42,ln=c,lw=1.5)
    tf=tf_box(s,fx,4.55,w,0.42,anchor=MSO_ANCHOR.MIDDLE,wrap=False); para(tf,[(t,dict(size=10,bold=True,color=c,name=MONO))],align=PP_ALIGN.CENTER)
    fx+=w
    if arrow:
        tf=tf_box(s,fx,4.55,0.5,0.42,anchor=MSO_ANCHOR.MIDDLE); para(tf,[("──▶",dict(size=11,color=FAINT))],align=PP_ALIGN.CENTER); fx+=0.5
chips(s,[{'t':"VS Code extension + CLI",'runs':[("VS Code extension ",{}),("+ CLI",dict(bold=True,color=BLACK))]},
         {'t':"Powered by Anthropic Claude",'runs':[("Powered by ",{}),("Anthropic Claude",dict(bold=True,color=BLACK))]},
         {'t':"63 automated tests · all passing",'runs':[("63",dict(bold=True,color=BLACK)),(" automated tests · all passing",{})]}],5.45)
footer(s,1)

# 2 · AGENDA
s=slide(PANEL)
header(s,"Agenda",["What we'll cover"],2,tsize=30)
items=[("01","The problem","why manual migration fails"),("02","What the platform does","in → out"),
("03","How it works","the four moves"),("04","Architecture","the layers, at a glance"),
("05","The pipeline","10 stages, end to end"),("06","The AI agent team","who does what"),
("07","The user journey","right-click → project"),("08","Self-healing verification","why you can trust it"),
("09","Grounding & safety","no hallucinations"),("10","Evidence, status, security","real runs · shipped scope"),
("11","Business case & FAQ","value, cost, questions"),("12","Roadmap & the ask","what's next")]
lx,_=cols(2,0.6); colw=(CW-0.6)/2
for i,(n,t,sub) in enumerate(items):
    x=lx[i%2]; y=BODY+0.2+(i//2)*0.72
    tf=tf_box(s,x,y,0.4,0.4); para(tf,[(n,dict(size=11,bold=True,color=GDEEP,name=MONO))])
    tf=tf_box(s,x+0.5,y,colw-2.4,0.4); para(tf,[(t,dict(size=13,bold=True,color=BLACK))])
    tf=tf_box(s,x+colw-2.3,y,2.3,0.4); para(tf,[(sub,dict(size=9.5,color=FAINT))],align=PP_ALIGN.RIGHT)
    rect(s,x,y+0.55,colw,0.012,fill=LINE)

# 3 · PROBLEM
s=slide()
header(s,"01 · The Problem",["Replatforming by hand is slow, expensive —","and quietly loses business logic"],3,tsize=27)
lx,cw=cols(3,0.32)
data=[("Months → Years","A mid-size Hybris estate holds hundreds of Java classes, a custom data model, live data, and scheduled jobs. Manual rewrites are measured in quarters."),
("Two-platform experts","It needs engineers senior in both Hybris and Salesforce — among the rarest, costliest skill combinations in the market."),
("Silent logic loss","Rules like “never accept a zero-value order” vanish in manual translation. Nobody notices — until production.")]
for i,(big,body) in enumerate(data):
    x=lx[i]; card(s,x,BODY+0.15,cw,3.9,top=RISK)
    tf=tf_box(s,x+0.24,BODY+0.42,cw-0.48,1.0); para(tf,[(big,dict(size=19,bold=True,color=RISK))],line=1.05)
    tf=tf_box(s,x+0.24,BODY+1.45,cw-0.48,2.3); para(tf,body,size=11,color=GRAY,line=1.28)
strip(s,"In plain words",[("Rewriting an old system by hand is like retyping a 1,000-page contract from memory — slow, and the fine print is what gets lost.",{})])

# 4 · WHAT IT DOES
s=slide(PANEL)
header(s,"02 · What It Does",["Point it at the old code. Get a working Salesforce project back"],4,tsize=25)
boxw=4.35; gap=CW-2*boxw
# IN
card(s,MX,BODY+0.2,boxw,3.7,top=COPPER)
tf=tf_box(s,MX+0.28,BODY+0.5,boxw-0.5,0.4); para(tf,[("IN — a Hybris codebase",dict(size=13.5,bold=True,color=BLACK))])
tf=tf_box(s,MX+0.28,BODY+1.05,boxw-0.5,2.5)
for it in ["Java business logic (orders, customers, pricing…)","The data model (items.xml)","Actual data records (ImpEx files)","Scheduled jobs (cronjobs + triggers)"]:
    para(tf,[("—  ",dict(color=COPPER)),(it,dict(color=GRAY))],size=11.5,line=1.2,sa=9)
# gear
gx=MX+boxw+gap/2-0.5
rect(s,gx,BODY+1.4,1.0,1.0,fill=BLACK,shape=MSO_SHAPE.OVAL)
tf=tf_box(s,gx,BODY+1.4,1.0,1.0,anchor=MSO_ANCHOR.MIDDLE); para(tf,[("AI AGENT\nTEAM",dict(size=8,bold=True,color=WHITE,name=MONO))],align=PP_ALIGN.CENTER,line=1.2)
rect(s,gx+0.82,BODY+2.12,0.12,0.12,fill=GREEN,shape=MSO_SHAPE.OVAL)
tf=tf_box(s,gx-0.35,BODY+2.55,1.7,0.3); para(tf,[("MINUTES–HOURS",dict(size=8.5,color=FAINT,name=MONO))],align=PP_ALIGN.CENTER)
# OUT
ox=MX+boxw+gap
card(s,ox,BODY+0.2,boxw,3.7,top=TEAL)
tf=tf_box(s,ox+0.28,BODY+0.5,boxw-0.5,0.4); para(tf,[("OUT — deployable Salesforce",dict(size=13.5,bold=True,color=BLACK))])
tf=tf_box(s,ox+0.28,BODY+1.05,boxw-0.5,2.5)
for it in ["Apex code in Salesforce enterprise patterns","A test class for every class","Data model, data & schedules — migrated","Confidence report: what to trust, what to review"]:
    para(tf,[("—  ",dict(color=TEAL)),(it,dict(color=GRAY))],size=11.5,line=1.2,sa=9)
strip(s,"In plain words",[("Old system in, new system out — with a quality report on top, like a surveyor's certificate handed over with the keys.",{})])

# 5 · HOW IT WORKS
s=slide()
header(s,"03 · How It Works",["Four moves — the way a strong engineering team would do it"],5,tsize=26)
lx,cw=cols(4,0.3)
rect(s,MX+cw/2,BODY+0.55,CW-cw,0.02,fill=LINE)
moves=[("1","Understand","Reads every file, maps dependencies, and writes down the business rules it finds — before writing any code."),
("2","Plan","Decides what each piece becomes — including what should not be custom code at all."),
("3","Build & Review","Writes the code and tests; a second AI reviews every piece skeptically before it's accepted."),
("4","Prove","Deploys to a real Salesforce environment and fixes real errors itself — in a loop, until green.")]
for i,(n,t,b) in enumerate(moves):
    x=lx[i]; cx=x+cw/2; last=(i==3)
    rc = GREEN if last else BLACK
    rect(s,cx-0.42,BODY+0.15,0.84,0.84,fill=WHITE,ln=rc,lw=2.0,shape=MSO_SHAPE.OVAL)
    tf=tf_box(s,cx-0.42,BODY+0.15,0.84,0.84,anchor=MSO_ANCHOR.MIDDLE); para(tf,[(n,dict(size=20,bold=True,color=(GDEEP if last else BLACK)))],align=PP_ALIGN.CENTER)
    tf=tf_box(s,x,BODY+1.2,cw,0.4); para(tf,[(t,dict(size=13,bold=True,color=BLACK))],align=PP_ALIGN.CENTER)
    tf=tf_box(s,x+0.05,BODY+1.65,cw-0.1,1.8); para(tf,b,size=10.5,color=GRAY,align=PP_ALIGN.CENTER,line=1.28)
strip(s,"In plain words",[("Read first, decide second, build with a reviewer looking over your shoulder, and don't call it done until it actually runs. Every decision is written down — no black box.",{})])

# 6 · ARCHITECTURE
s=slide(PANEL)
header(s,"04 · Architecture",["Four layers, one engine"],6,tsize=30)
layers=[("INTERFACES",BLACK,["VS Code extension — right-click","Command line — for CI/CD","both drive the same engine"]),
("ORCHESTRATION",GDEEP,["Agentic mode — Planner · Builder · Critic · Verifier","Linear mode — fixed 10-stage pipeline"]),
("SHARED STAGE FUNCTIONS",BLACK,["parse Java","derive schema","generate Apex","validate & repair","data · jobs","verify · report"]),
("AI PROVIDERS",BLACK,["Anthropic Claude — best quality","OpenRouter — cheap iteration","Mock — free, offline, zero exposure"])]
ly=BODY+0.05; lh=0.86
for i,(name,lc,nodes) in enumerate(layers):
    rect(s,MX,ly,CW,lh,fill=WHITE,ln=LINE,lw=0.75)
    rect(s,MX,ly,1.95,lh,fill=lc)
    tf=tf_box(s,MX+0.16,ly,1.7,lh,anchor=MSO_ANCHOR.MIDDLE); para(tf,[(name,dict(size=8.5,bold=True,color=WHITE,name=MONO,track=0.8))],line=1.1)
    nx=MX+2.15
    for nd in nodes:
        w=0.28+len(nd)*0.066
        rect(s,nx,ly+lh/2-0.19,w,0.38,fill=PANEL,ln=LINE,lw=0.5)
        tf=tf_box(s,nx+0.1,ly+lh/2-0.19,w-0.2,0.38,anchor=MSO_ANCHOR.MIDDLE); para(tf,[(nd,dict(size=9.5,bold=True,color=BLACK))])
        nx+=w+0.14
    if i<3:
        tf=tf_box(s,MX,ly+lh-0.02,CW,0.24); para(tf,[("▼",dict(size=9,color=FAINT))],align=PP_ALIGN.CENTER)
    ly+=lh+0.24
strip(s,"In plain words",[("Two “driving modes” — a smart agent team, or a simple assembly line — share the same proven tools underneath, so every fix lands in both. And the AI brain is swappable, including a free offline one.",{})])

# 7 · PIPELINE
s=slide()
header(s,"05 · The Pipeline",["Ten stages, end to end"],7,tsize=30)
stg=[("01","Crawl & schedule","Group classes by domain; sort so dependencies convert first.",0),
("02","Call graph","Map who-calls-whom for the visual dashboard.",0),
("03","Ingest","Parse the Java and data-model definitions precisely.",0),
("04","Derive schema","Build the target object/field catalog — the source of truth.",0),
("05","Comprehend","AI summarises each class: purpose, queries, rules.",1),
("06","Generate","AI writes the Apex + tests, grounded in the schema.",1),
("07","Validate & repair","Safety checks; failures loop back to the AI to fix.",1),
("08","Reconcile & metadata","Fill schema gaps with evidence; emit objects & fields.",0),
("09","Data & jobs","ImpEx → CSVs; cron triggers → scheduling runbook.",0),
("10","Verify & report","Real-org deploy + self-heal; confidence scores.",1)]
lx,cw=cols(5,0.22); ch=1.55
for i,(n,t,b,ai) in enumerate(stg):
    x=lx[i%5]; y=BODY+0.05+(i//5)*(ch+0.24)
    card(s,x,y,cw,ch,top=None)
    tf=tf_box(s,x+0.16,y+0.14,cw-0.3,0.24); para(tf,[(n,dict(size=9,bold=True,color=GDEEP,name=MONO))])
    if ai:
        rect(s,x+cw-0.5,y+0.13,0.36,0.24,fill=BLACK)
        tf=tf_box(s,x+cw-0.5,y+0.12,0.36,0.24,anchor=MSO_ANCHOR.MIDDLE); para(tf,[("AI",dict(size=8,bold=True,color=WHITE,name=MONO))],align=PP_ALIGN.CENTER)
    tf=tf_box(s,x+0.16,y+0.4,cw-0.3,0.35); para(tf,[(t,dict(size=11,bold=True,color=BLACK))],line=1.05)
    tf=tf_box(s,x+0.16,y+0.78,cw-0.3,0.7); para(tf,b,size=8.6,color=GRAY,line=1.2)
strip(s,"In plain words",[("The ",{}),("AI",dict(bold=True,color=BLACK)),(" is used only where judgement is needed — understanding and writing code. Everything mechanical is ordinary tested software: fast, free, repeatable.",{})])

# 8 · AGENTIC CORE
s=slide(PANEL)
header(s,"06 · The Agentic Core",["A manager, a shared whiteboard, and four specialists"],8,tsize=27)
rect(s,MX,BODY+0.2,CW,0.62,fill=BLACK)
tf=tf_box(s,MX,BODY+0.3,CW,0.44,anchor=MSO_ANCHOR.MIDDLE)
para(tf,[("Orchestrator   ",dict(size=13,bold=True,color=WHITE)),("ROUTES WORK · LOOPS UNTIL EVERYTHING IS VERIFIED",dict(size=8.5,color=GREEN,name=MONO,track=1.0))],align=PP_ALIGN.CENTER)
tf=tf_box(s,MX,BODY+0.85,CW,0.22); para(tf,[("▲   ▼",dict(size=9,color=FAINT))],align=PP_ALIGN.CENTER)
rect(s,MX,BODY+1.1,CW,0.66,ln=GREEN,lw=1.5)
tf=tf_box(s,MX,BODY+1.1,CW,0.66,anchor=MSO_ANCHOR.MIDDLE)
para(tf,[("The Blackboard",dict(size=12,bold=True,color=GDEEP)),(" — shared state:  schema · migration plan · generated artifacts · ",dict(size=11,color=INK)),("decision log",dict(size=11,bold=True,color=GDEEP)),(" · open questions",dict(size=11,color=INK))],align=PP_ALIGN.CENTER)
tf=tf_box(s,MX,BODY+1.8,CW,0.22); para(tf,[("▲   ▼",dict(size=9,color=FAINT))],align=PP_ALIGN.CENTER)
lx,cw=cols(4,0.22)
ag=[("STRATEGY","Planner"),("EXECUTION","Builder"),("QUALITY GATE","Critic"),("PROOF","Verifier")]
for i,(role,nm) in enumerate(ag):
    x=lx[i]; card(s,x,BODY+2.1,cw,0.85,top=None)
    tf=tf_box(s,x,BODY+2.25,cw,0.28); para(tf,[(role,dict(size=8,bold=True,color=GDEEP,name=MONO,track=0.6))],align=PP_ALIGN.CENTER)
    tf=tf_box(s,x,BODY+2.5,cw,0.35); para(tf,[(nm,dict(size=14,bold=True,color=BLACK))],align=PP_ALIGN.CENTER)
strip(s,"In plain words",[("Four specialists around a shared whiteboard — so work can be sent back (Critic finds a problem → Builder fixes it), exactly like a real team. The whiteboard becomes a readable document after every run.",{})])

# 9 · FOUR AGENTS
s=slide()
header(s,"06 · The Agent Roles",["Who does what — and what each one replaces"],9,tsize=28)
lx,cw=cols(4,0.24)
ag=[("STRATEGY","The Planner","Routes every piece: build as Apex, recommend a native Salesforce product, or skip — with a written rationale.","“Pricing rules? That's Salesforce CPQ — don't hand-build it.”","blind, hard-coded translation of everything."),
("EXECUTION","The Builder","Writes the Apex + a test class per class, grounded in the real schema and a built-in best-practice base.","“Bulk-safe, secure, in the house pattern — like a senior dev.”","months of manual line-by-line rewriting."),
("QUALITY GATE","The Critic","Adversarially reviews every artifact: behavior preserved? Secure? Real problems trigger a fix-and-re-review.","“The zero-total rule got dropped on one path — blocked.”","hoping a human catches everything later."),
("PROOF","The Verifier","Deploys to a real org, reads real compiler errors, heals them, and re-deploys until green.","“Not 'looks right' — it actually ran.”","“trust me, it compiles.”")]
for i,(role,nm,body,q,repl) in enumerate(ag):
    x=lx[i]; card(s,x,BODY+0.15,cw,4.9,top=GREEN)
    tf=tf_box(s,x+0.2,BODY+0.42,cw-0.4,0.24); para(tf,[(role,dict(size=8,bold=True,color=GDEEP,name=MONO,track=0.8))])
    tf=tf_box(s,x+0.2,BODY+0.7,cw-0.4,0.35); para(tf,[(nm,dict(size=14.5,bold=True,color=BLACK))])
    tf=tf_box(s,x+0.2,BODY+1.15,cw-0.4,1.8); para(tf,body,size=10.3,color=GRAY,line=1.26)
    rect(s,x+0.2,BODY+2.75,0.04,0.75,fill=GREEN)
    tf=tf_box(s,x+0.34,BODY+2.75,cw-0.54,0.9); para(tf,[(q,dict(italic=True,color=INK))],size=9.8,line=1.22)
    rect(s,x+0.2,BODY+3.75,cw-0.4,1.0,fill=PANEL)
    tf=tf_box(s,x+0.34,BODY+3.87,cw-0.6,0.85); para(tf,[("REPLACES — ",dict(size=8,bold=True,color=COPPER,name=MONO)),(repl,dict(size=9.2,color=GRAY))],line=1.2)

# 10 · USER JOURNEY
s=slide(PANEL)
header(s,"07 · The User Journey",["From right-click to finished project"],10,tsize=28)
lx,cw=cols(5,0.22)
js=[("1","Configure once","Pick the AI provider and paste your key — or choose free mock mode."),
("2","Right-click","On the Hybris folder → “H2A: Migrate to Apex.” First run self-installs."),
("3","Watch","A dashboard streams progress: domains, plan decisions, review results."),
("4","Receive","A salesforce_<project> folder appears — a complete, deployable project."),
("5","Review & ship","Open the report; review what's flagged; deploy with one command.")]
for i,(n,t,b) in enumerate(js):
    x=lx[i]; card(s,x,BODY+0.2,cw,2.5,top=None)
    dc = GDEEP if i==4 else BLACK
    rect(s,x+0.2,BODY+0.42,0.4,0.4,fill=dc,shape=MSO_SHAPE.OVAL)
    tf=tf_box(s,x+0.2,BODY+0.42,0.4,0.4,anchor=MSO_ANCHOR.MIDDLE); para(tf,[(n,dict(size=11,bold=True,color=WHITE))],align=PP_ALIGN.CENTER)
    tf=tf_box(s,x+0.2,BODY+0.95,cw-0.4,0.35); para(tf,[(t,dict(size=11.5,bold=True,color=BLACK))])
    tf=tf_box(s,x+0.2,BODY+1.35,cw-0.4,1.1); para(tf,b,size=9.5,color=GRAY,line=1.24)
rect(s,MX,BODY+2.95,CW,0.6,fill=PANEL,ln=LINE,lw=0.5)
tf=tf_box(s,MX+0.25,BODY+2.95,CW-0.5,0.6,anchor=MSO_ANCHOR.MIDDLE)
para(tf,[("BEHIND THE SCENES   ",dict(size=8.5,bold=True,color=GDEEP,name=MONO)),("the extension bundles the entire Python engine — no manual installation; the same run is scriptable from a terminal for CI/CD.",dict(size=10.5,color=GRAY))])
strip(s,"In plain words",[("For the person using it, the whole platform is one right-click and one folder of results.",{})])

# 11 · SELF-HEALING
s=slide()
header(s,"08 · Self-Healing Verification",["It doesn't stop at “the AI wrote code.” It proves the code runs"],11,tsize=24)
lx,cw=cols(4,0.22)
loop=[("Deploy for real","Validation-only deploy to an actual Salesforce org"),
("Read real errors","The actual compiler output — not a guess"),
("Fix automatically","Three healing modes (below)"),
("Repeat until green","Anything unresolved is flagged for a human")]
for i,(t,b) in enumerate(loop):
    x=lx[i]; card(s,x,BODY+0.1,cw,1.15,top=None)
    tf=tf_box(s,x+0.16,BODY+0.28,cw-0.3,0.3); para(tf,[(t,dict(size=11.5,bold=True,color=BLACK))],align=PP_ALIGN.CENTER)
    tf=tf_box(s,x+0.16,BODY+0.62,cw-0.3,0.5); para(tf,b,size=9,color=GRAY,align=PP_ALIGN.CENTER,line=1.2)
    if i<3:
        tf=tf_box(s,x+cw,BODY+0.45,0.22,0.4,anchor=MSO_ANCHOR.MIDDLE); para(tf,[("→",dict(size=13,color=FAINT))],align=PP_ALIGN.CENTER)
tf=tf_box(s,MX,BODY+1.4,CW,0.3); para(tf,[("↺  BOUNDED LOOP — RUNS UNTIL GREEN, NEVER SILENTLY SHIPS",dict(size=9,bold=True,color=GDEEP,name=MONO,track=1.0))],align=PP_ALIGN.CENTER)
lx3,cw3=cols(3,0.3)
heal=[("Metadata healing","“Missing field” error? If the field is genuinely used in the original Java — proven by evidence — it's added to the data model. Never guessed."),
("Code repair","Compile errors are fed back to the AI with the real error message, and the class is rewritten and re-tried."),
("Coverage healing","Salesforce requires 75% test coverage to deploy. Below it? The tool writes more tests — error paths, bulk scenarios — until it clears the bar.")]
for i,(t,b) in enumerate(heal):
    x=lx3[i]; card(s,x,BODY+1.85,cw3,2.55,top=BLACK)
    tf=tf_box(s,x+0.22,BODY+2.1,cw3-0.44,0.35); para(tf,[(t,dict(size=13,bold=True,color=BLACK))])
    tf=tf_box(s,x+0.22,BODY+2.55,cw3-0.44,1.7); para(tf,b,size=10.5,color=GRAY,line=1.28)

# 12 · GROUNDING
s=slide(PANEL)
header(s,"09 · Grounding & Safety",["Three mechanisms that stop the AI from making things up"],12,tsize=26)
lx,cw=cols(3,0.32)
g=[("1 · Schema grounding","Before writing a single query, the AI is shown the exact catalog of objects and fields from your real data model — and everything it writes is checked against that catalog afterwards. Anything that doesn't exist is caught."),
("2 · Evidence-based reconciliation","Generated code references a field the model never declared? The system checks the original Java source. Genuinely used → added, properly typed. No evidence → flagged for human review. Never silently guessed."),
("3 · A built-in knowledge base","Salesforce's governor limits, security rules, and enterprise patterns are bundled as a reference library the agents retrieve and cite while writing and reviewing — facts, not memory.")]
for i,(t,b) in enumerate(g):
    x=lx[i]; card(s,x,BODY+0.15,cw,3.9,top=GREEN)
    tf=tf_box(s,x+0.24,BODY+0.42,cw-0.48,0.35); para(tf,[(t,dict(size=13.5,bold=True,color=BLACK))])
    tf=tf_box(s,x+0.24,BODY+0.9,cw-0.48,2.9); para(tf,b,size=11,color=GRAY,line=1.3)
strip(s,"In plain words",[("The AI works “open book”: it can only use what provably exists, its claims are checked against the source, and it looks rules up instead of recalling them.",{})])

# 13 · DELIVERABLE
s=slide()
header(s,"10 · The Deliverable",["What lands in the output folder"],13,tsize=30)
lx,cw=cols(3,0.3)
inv=[("force-app/…/classes","Apex code + tests","Selectors, Services, REST controllers, scheduled jobs — each with its own test class."),
("force-app/…/objects","Data model","Custom objects, fields, picklists, relationships, required/unique constraints."),
("data/ · DATA_MIGRATION.md","The data","Load-ready CSVs plus a safe, re-runnable import runbook."),
("CRON_JOBS.md · schedule.apex","Schedules","Every timed job translated, same timing, with a ready-to-run script."),
("MIGRATION_PLAN.md","The decisions","The Planner's call on every class, the Critic's findings, the full decision log."),
("FEASIBILITY_REPORT.md","The scorecard","Validation results, a High/Medium/Low confidence score per class, deploy status, cost.")]
for i,(code,t,b) in enumerate(inv):
    x=lx[i%3]; y=BODY+0.15+(i//3)*1.9
    card(s,x,y,cw,1.7,top=None)
    rect(s,x+0.2,y+0.2,min(cw-0.4,0.2+len(code)*0.062),0.32,fill=PANEL)
    tf=tf_box(s,x+0.28,y+0.2,cw-0.5,0.32,anchor=MSO_ANCHOR.MIDDLE); para(tf,[(code,dict(size=9,color=GDEEP,name=MONO))])
    tf=tf_box(s,x+0.2,y+0.58,cw-0.4,0.3); para(tf,[(t,dict(size=12,bold=True,color=BLACK))])
    tf=tf_box(s,x+0.2,y+0.92,cw-0.4,0.7); para(tf,b,size=9.8,color=GRAY,line=1.24)
strip(s,"In plain words",[("Not just code — a complete package: the system, its data, its schedules, and the paper trail that tells your reviewers exactly where to look.",{})])

# 14 · EVIDENCE
s=slide(PANEL)
header(s,"11 · Evidence",["Two unscripted moments from a real run"],14,tsize=30)
boxw=(CW-0.35)/2
cases=[("THE PLANNER'S JUDGEMENT",TEAL,"It refused to write unnecessary code.","Handed a custom discount/promotions engine, the Planner didn't translate it — it recommended Salesforce CPQ, the native product built for pricing rules, and generated no custom code for it.","Its own words: “Discount/promo-code pricing rules are a textbook fit for Salesforce CPQ.” Architect-level judgement — less code to own, less debt."),
("THE CRITIC'S CATCH",GDEEP,"It caught a silently-lost business rule.","The original Java rejected non-positive order totals. The first translation compiled perfectly — but dropped that check on one path. The Critic caught it and blocked the class as “needs review.”","This is exactly the bug class that slips through manual migrations and surfaces in production. Here it never reached a human unflagged.")]
for i,(tag,tc,h,body,verdict) in enumerate(cases):
    x=MX+i*(boxw+0.35); card(s,x,BODY+0.15,boxw,4.4,top=BLACK)
    rect(s,x+0.26,BODY+0.42,0.28+len(tag)*0.058,0.32,fill=(PANEL if i==1 else RGBColor(0xe6,0xf1,0xf6)))
    tf=tf_box(s,x+0.36,BODY+0.42,boxw-0.5,0.32,anchor=MSO_ANCHOR.MIDDLE); para(tf,[(tag,dict(size=8.5,bold=True,color=tc,name=MONO,track=0.8))])
    tf=tf_box(s,x+0.26,BODY+0.9,boxw-0.5,0.4); para(tf,[(h,dict(size=15,bold=True,color=BLACK))])
    tf=tf_box(s,x+0.26,BODY+1.45,boxw-0.5,1.6); para(tf,body,size=11,color=GRAY,line=1.28)
    rect(s,x+0.26,BODY+3.05,boxw-0.52,1.1,fill=PANEL)
    tf=tf_box(s,x+0.42,BODY+3.2,boxw-0.8,0.9); para(tf,verdict,size=10.3,color=INK,line=1.26)

# 15 · STATUS
s=slide()
header(s,"12 · Current Implementation",["Shipped and working today"],15,tsize=30)
rows=[("Code migration","Java business logic → best-practice Apex, one test class per class","SHIPPED"),
("Data model","Objects, fields, picklists, relationships, constraints — from items.xml","SHIPPED"),
("Data records","ImpEx → load-ready CSVs + re-runnable import runbook","SHIPPED"),
("Scheduled jobs","Cronjobs → Salesforce scheduled Apex, identical timing","SHIPPED"),
("Agent team + self-healing verification","Planner · Builder · Critic · Verifier; deploy loop; confidence scoring","SHIPPED"),
("Workflows & storefront APIs","Hybris processes → Flow; OCC REST → Apex REST","ROADMAP")]
ty=BODY+0.15; rh=0.62
# header row
tf=tf_box(s,MX+0.1,ty,4,0.3); para(tf,[("CAPABILITY",dict(size=8.5,color=GRAY,name=MONO,track=1.0))])
tf=tf_box(s,MX+4.3,ty,5,0.3); para(tf,[("WHAT IT COVERS",dict(size=8.5,color=GRAY,name=MONO,track=1.0))])
tf=tf_box(s,PW-MX-1.3,ty,1.3,0.3); para(tf,[("STATUS",dict(size=8.5,color=GRAY,name=MONO,track=1.0))])
rect(s,MX,ty+0.32,CW,0.022,fill=BLACK)
for i,(cap,cov,st) in enumerate(rows):
    y=ty+0.4+i*rh
    tf=tf_box(s,MX+0.1,y,4.1,rh-0.1,anchor=MSO_ANCHOR.MIDDLE); para(tf,[(cap,dict(size=11,bold=True,color=BLACK))],line=1.1)
    tf=tf_box(s,MX+4.3,y,5.2,rh-0.1,anchor=MSO_ANCHOR.MIDDLE); para(tf,[(cov,dict(size=10.5,color=GRAY))],line=1.15)
    done = st=="SHIPPED"
    bw=0.98 if done else 1.02
    rect(s,PW-MX-bw,y+rh/2-0.16,bw,0.32,fill=(RGBColor(0xea,0xf3,0xda) if done else PANEL))
    tf=tf_box(s,PW-MX-bw,y+rh/2-0.16,bw,0.32,anchor=MSO_ANCHOR.MIDDLE); para(tf,[(st,dict(size=8.5,bold=True,color=(GDEEP if done else FAINT),name=MONO))],align=PP_ALIGN.CENTER)
    rect(s,MX,y+rh-0.06,CW,0.012,fill=LINE)
chips(s,[{'t':"63 automated tests, all passing",'runs':[("63",dict(bold=True,color=BLACK)),(" automated tests, all passing",{})]},
         {'t':"3 AI providers incl. a free offline mode",'runs':[("3",dict(bold=True,color=BLACK)),(" AI providers incl. a free offline mode",{})]},
         {'t':"every run reports its own cost",'runs':[("every run reports its own ",{}),("cost",dict(bold=True,color=BLACK))]}],ty+0.4+6*rh+0.15)

# 16 · SECURITY
s=slide(PANEL)
header(s,"13 · Security & Data Handling",["Your code, your keys, your choice of where it goes"],16,tsize=27)
lx,cw=cols(2,0.35)
sec=[("1","You control the AI provider","Source code goes only to the vendor you configure. Keys live in your settings — never in source control, never in the shipped product (audited every release)."),
("2","A zero-exposure mode","Mock mode runs the entire pipeline with nothing leaving the machine — for sensitive codebases and free evaluation."),
("3","Secure code by default","Generated Apex enforces Salesforce field-level security and record-sharing as a house standard, not an afterthought."),
("4","Nothing destructive, ever","Deployments are validation-only dry runs. A human holds the final go-live button. Every AI decision is logged and auditable.")]
for i,(n,t,b) in enumerate(sec):
    x=lx[i%2]; y=BODY+0.15+(i//2)*1.6
    card(s,x,y,cw,1.42,top=None)
    rect(s,x+0.22,y+0.24,0.42,0.42,fill=BLACK,shape=MSO_SHAPE.OVAL)
    tf=tf_box(s,x+0.22,y+0.24,0.42,0.42,anchor=MSO_ANCHOR.MIDDLE); para(tf,[(n,dict(size=12,bold=True,color=GREEN))],align=PP_ALIGN.CENTER)
    tf=tf_box(s,x+0.82,y+0.22,cw-1.05,0.32); para(tf,[(t,dict(size=12.5,bold=True,color=BLACK))])
    tf=tf_box(s,x+0.82,y+0.56,cw-1.05,0.8); para(tf,b,size=10,color=GRAY,line=1.24)
strip(s,"On the roadmap",[("Private/VPC model hosting and approval-gated review workflows for regulated environments.",{})])

# 17 · BUSINESS CASE
s=slide()
header(s,"14 · The Business Case",["What changes, concretely"],17,tsize=30)
ty=BODY+0.2; rh=0.72
tf=tf_box(s,MX+0.1,ty,5,0.3); para(tf,[("WITHOUT THE PLATFORM",dict(size=8.5,color=GRAY,name=MONO,track=1.0))])
tf=tf_box(s,MX+CW/2+0.1,ty,5,0.3); para(tf,[("WITH THE PLATFORM",dict(size=8.5,color=GRAY,name=MONO,track=1.0))])
rect(s,MX,ty+0.32,CW,0.022,fill=BLACK)
biz=[("Months of manual rewriting","A working first draft in minutes–hours"),
("Business rules silently lost","An AI reviewer checks rule preservation; a parity score proves it"),
("“Trust me, it compiles”","Deployed to a real org and self-corrected until verifiably green"),
("Everything becomes custom code to own forever","Native Salesforce products recommended where they fit — less debt"),
("A black-box engagement","A decision log and a confidence score on every class")]
for i,(a,b) in enumerate(biz):
    y=ty+0.4+i*rh
    tf=tf_box(s,MX+0.1,y,CW/2-0.3,rh-0.12,anchor=MSO_ANCHOR.MIDDLE); para(tf,[(a,dict(size=11.5,bold=True,color=RISK))],line=1.12)
    tf=tf_box(s,MX+CW/2+0.1,y,CW/2-0.3,rh-0.12,anchor=MSO_ANCHOR.MIDDLE); para(tf,[(b,dict(size=11.5,bold=True,color=GDEEP))],line=1.12)
    rect(s,MX,y+rh-0.06,CW,0.012,fill=LINE)
strip(s,"The honest positioning",[("A strong, verified first draft plus a prioritized review list. It makes expert reviewers dramatically faster — it doesn't remove them, and we don't pretend it does.",{})])

# 18 · FAQ
s=slide(PANEL)
header(s,"15 · Common Questions",["The questions every room asks"],18,tsize=30)
lx,cw=cols(2,0.32)
faq=[("Does this replace our developers?","No — it replaces the grind. Experts review a scored draft instead of rewriting; confidence scores show exactly where to look."),
("What if the AI makes something up?","Three nets: it can only use fields that provably exist; a second AI reviews every class; the code must compile against a real org."),
("Is our source code safe?","It goes only to the provider you configure — or nowhere at all in mock mode. No secrets are ever bundled; deploys are dry-run only."),
("What does a run cost?","Every report itemizes its own AI usage. Caching keeps repeat costs low; cheaper models handle the simple steps automatically."),
("What can't it do yet?","Hybris workflow processes and storefront REST APIs are next on the roadmap — and it tells you honestly what it skipped, and why."),
("Can we try it on our codebase?","Yes — today. Free mock mode first (zero exposure), then a real scored run. It's a right-click in VS Code.")]
for i,(q,a) in enumerate(faq):
    x=lx[i%2]; y=BODY+0.15+(i//2)*1.5
    card(s,x,y,cw,1.32,top=None)
    tf=tf_box(s,x+0.22,y+0.2,cw-0.44,0.32); para(tf,[("Q.  ",dict(size=11,bold=True,color=GDEEP,name=MONO)),(q,dict(size=11.5,bold=True,color=BLACK))])
    tf=tf_box(s,x+0.22,y+0.6,cw-0.44,0.7); para(tf,a,size=10,color=GRAY,line=1.26)

# 19 · ROADMAP
s=slide()
header(s,"16 · Where This Is Going",["Built in phases — each one de-risks the next"],19,tsize=28)
lx,cw=cols(5,0.2)
ph=[("PHASE 0 ✓",GREEN,GDEEP,"Prove correctness","Self-healing deploys, confidence scoring, behavior-parity checks"),
("PHASE 1 ✓",GREEN,GDEEP,"The agent team","Planner, Builder, Critic, Verifier over a shared whiteboard"),
("PHASE 2 ◐",WARN,WARN,"The whole platform","Data ✓ · schema ✓ · jobs ✓ · workflows & storefront APIs next"),
("PHASE 3",LINE,FAINT,"Enterprise-grade","Review workspace, audit trail, private/VPC models, org-aware reuse"),
("PHASE 4",LINE,FAINT,"It learns","Every migration makes the next better and cheaper; assessment product")]
for i,(n,topc,nc,t,b) in enumerate(ph):
    x=lx[i]; card(s,x,BODY+0.3,cw,3.0,top=topc)
    tf=tf_box(s,x+0.16,BODY+0.55,cw-0.3,0.28); para(tf,[(n,dict(size=8.5,bold=True,color=nc,name=MONO,track=0.6))])
    tf=tf_box(s,x+0.16,BODY+0.85,cw-0.3,0.6); para(tf,[(t,dict(size=11.5,bold=True,color=BLACK))],line=1.1)
    tf=tf_box(s,x+0.16,BODY+1.45,cw-0.3,1.4); para(tf,b,size=9.3,color=GRAY,line=1.28)
strip(s,"The through-line",[("Every phase moves along one axis: from output you have to trust → to output you can verify. That axis is the moat.",{})])

# 20 · CLOSE
s=slide(BLACK)
rect(s,4.6,6.7,2.5,2.5,ln=RGBColor(0x2a,0x2a,0x2a),lw=1.0,shape=MSO_SHAPE.OVAL)
tf=tf_box(s,0,1.7,PW,0.4); para(tf,[("● THE ASK",dict(size=10,bold=True,color=GREEN,name=MONO,track=2.0))],align=PP_ALIGN.CENTER)
tf=tf_box(s,1.0,2.5,PW-2,1.6)
para(tf,[("Give us one real slice of a Hybris codebase.",dict(size=27,bold=True,color=WHITE))],align=PP_ALIGN.CENTER,line=1.12,sa=4)
para(tf,[("We'll hand back verified Salesforce — and the receipts",dict(size=27,bold=True,color=WHITE)),(" ●",dict(size=27,bold=True,color=GREEN))],align=PP_ALIGN.CENTER,line=1.12)
tf=tf_box(s,2.2,4.5,PW-4.4,0.9); para(tf,"A pilot costs days, not months. Free mock mode first, a scored real run second, your team's verdict third.",size=13,color=RGBColor(0xa7,0xa8,0xaa),align=PP_ALIGN.CENTER,line=1.3)
# chips centered
citems=[("Live demo today",[("Live demo ",dict(color=RGBColor(0xa7,0xa8,0xaa))),("today",dict(bold=True,color=WHITE))]),
("Full docs: PRD · TDD · flows",[("Full docs: ",dict(color=RGBColor(0xa7,0xa8,0xaa))),("PRD · TDD · flows · security",dict(bold=True,color=WHITE))]),
("Pilot-ready",[("Pilot-ready",dict(bold=True,color=WHITE))])]
cw_c=[0.42+len(t)*0.088 for t,_ in citems]
tot=sum(cw_c)+0.4
cx=(PW-tot)/2
for (t,runs),w in zip(citems,cw_c):
    rect(s,cx,5.6,w,0.42,fill=RGBColor(0x18,0x18,0x18),ln=RGBColor(0x3a,0x3a,0x3a),lw=0.75)
    tf=tf_box(s,cx+0.18,5.6,w-0.36,0.42,anchor=MSO_ANCHOR.MIDDLE,wrap=False); para(tf,runs,size=9,name=MONO,align=PP_ALIGN.CENTER)
    cx+=w+0.2
footer(s,20,dark=True)

import sys
out=sys.argv[1] if len(sys.argv)>1 else "DEMO_DECK.pptx"
prs.save(out)
print("saved",out,"slides:",len(prs.slides._sldIdLst))
