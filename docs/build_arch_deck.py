#!/usr/bin/env python
"""Build docs/ARCHITECTURE_DECK.pptx — the 10-slide H2A architecture deck (A4 landscape).

Shares the visual language of build_deck.py (same palette, geometry, and helpers) so the
two decks sit together without looking like they came from different projects.

Slide 3 embeds docs/architecture-diagram.png — run build_arch_image.py first if it is
missing or stale.

Regenerate with:  python docs/build_arch_deck.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DOCS = Path(__file__).resolve().parent

# ── palette (identical to build_deck.py, plus a purple for the model layer) ──
BLACK=RGBColor(0x0E,0x0E,0x0E); INK=RGBColor(0x22,0x24,0x27); GRAY=RGBColor(0x53,0x56,0x5A)
FAINT=RGBColor(0x8f,0x92,0x96); LINE=RGBColor(0xd2,0xd2,0xcf); PANEL=RGBColor(0xf3,0xf3,0xf1)
WHITE=RGBColor(0xFF,0xFF,0xFF); GREEN=RGBColor(0x86,0xBC,0x25); GDEEP=RGBColor(0x04,0x6a,0x38)
TEAL=RGBColor(0x00,0x7c,0xb0); COPPER=RGBColor(0xa9,0x62,0x2f); RISK=RGBColor(0xb3,0x28,0x2d)
PURPLE=RGBColor(0x6b,0x4f,0xa8)
SANS="Arial"; MONO="Courier New"
N_SLIDES=10

# ── geometry (A4 landscape) ──
PW,PH = 11.69, 8.27
MX = 0.66
CW = PW-2*MX
EY_Y=0.52; TITLE_Y=0.9; RULE_Y=1.82; BODY=2.08
STRIP_Y=7.06; STRIP_H=0.78; FOOT_Y=7.98

prs=Presentation(); prs.slide_width=Inches(PW); prs.slide_height=Inches(PH)

def slide(bg=WHITE):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.shadow.inherit=False; r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background()
    return s

def rect(s,l,t,w,h,fill=None,ln=None,lw=0.75,shape=MSO_SHAPE.RECTANGLE):
    sp=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h)); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if ln is None: sp.line.fill.background()
    else: sp.line.color.rgb=ln; sp.line.width=Pt(lw)
    return sp

def tf_box(s,l,t,w,h,anchor=MSO_ANCHOR.TOP,wrap=True):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return tf

def _track(run,pts): run._r.get_or_add_rPr().set('spc',str(int(pts*100)))

def para(tf, runs, align=PP_ALIGN.LEFT, size=11, color=INK, bold=False, name=SANS,
         italic=False, sa=3, sb=0, line=1.08):
    p0=tf.paragraphs[0]
    p = p0 if (not p0.runs) else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(sa); p.space_before=Pt(sb); p.line_spacing=line
    if isinstance(runs,str): runs=[(runs,{})]
    for txt,st in runs:
        r=p.add_run(); r.text=txt; f=r.font
        f.size=Pt(st.get('size',size)); f.bold=st.get('bold',bold); f.italic=st.get('italic',italic)
        f.name=st.get('name',name); f.color.rgb=st.get('color',color)
        if st.get('track'): _track(r,st['track'])
    return p

def eyebrow(s,kick,dark=False):
    rect(s,MX,EY_Y+0.02,0.13,0.13,fill=GREEN,shape=MSO_SHAPE.OVAL)
    tf=tf_box(s,MX+0.27,EY_Y-0.04,CW-0.3,0.32)
    para(tf,[(kick.upper(),dict(size=10.5,bold=True,color=(FAINT if dark else GRAY),name=MONO,track=2.0))])

def title(s,lines,size=29,dark=False,y=TITLE_Y):
    ink=WHITE if dark else BLACK
    tf=tf_box(s,MX,y,CW,1.05)
    for i,ln in enumerate(lines):
        runs=[(ln,dict(size=size,bold=True,color=ink,name=SANS))]
        if i==len(lines)-1: runs.append((".",dict(size=size,bold=True,color=GREEN,name=SANS)))
        para(tf,runs,size=size,line=1.02,sa=0)

def footer(s,page,dark=False):
    col=RGBColor(0x6a,0x6d,0x70) if dark else FAINT
    tf=tf_box(s,MX,FOOT_Y,4.6,0.25); para(tf,[("H2A MIGRATOR · ARCHITECTURE",dict(size=8,color=col,name=MONO,track=1.2))])
    tf2=tf_box(s,PW-MX-2.2,FOOT_Y,2.0,0.25); para(tf2,[(f"{page:02d} / {N_SLIDES}",dict(size=8,color=col,name=MONO,track=1.2))],align=PP_ALIGN.RIGHT)
    rect(s,PW-MX-0.14,FOOT_Y+0.025,0.09,0.09,fill=GREEN,shape=MSO_SHAPE.OVAL)

def header(s,kick,lines,page,tsize=29):
    eyebrow(s,kick); title(s,lines,size=tsize)
    rect(s,MX,RULE_Y,CW,0.02,fill=LINE)
    footer(s,page)

def card(s,l,t,w,h,top=BLACK,fill=WHITE,border=LINE):
    r=rect(s,l,t,w,h,fill=fill,ln=border,lw=0.75)
    if top: rect(s,l,t,w,0.055,fill=top)
    return r

def strip(s,label,body_runs,y=STRIP_Y,h=STRIP_H):
    rect(s,MX,y,CW,h,fill=PANEL)
    rect(s,MX,y,0.07,h,fill=GREEN)
    tf=tf_box(s,MX+0.28,y+0.12,CW-0.55,h-0.22)
    para(tf,[(label.upper(),dict(size=8.5,bold=True,color=GDEEP,name=MONO,track=1.4))],sa=4)
    para(tf,body_runs,size=11.5,color=INK,line=1.16)

def chips(s,items,y,l=MX):
    x=l
    for it in items:
        w=0.36+len(it['t'])*0.086
        rect(s,x,y,w,0.38,fill=WHITE,ln=LINE,lw=0.75)
        tf=tf_box(s,x+0.15,y,w-0.3,0.38,anchor=MSO_ANCHOR.MIDDLE,wrap=False)
        para(tf,it['runs'],size=9,name=MONO,color=GRAY,line=1.0)
        x+=w+0.2

def cols(n,gap,l=MX,w=CW):
    cw=(w-(n-1)*gap)/n
    return [l+i*(cw+gap) for i in range(n)], cw

def bullets(tf,items,dot=GREEN,size=11.5,line=1.24,sa=7,color=GRAY):
    for it in items:
        runs=[("—  ",dict(color=dot))]
        runs += it if isinstance(it,list) else [(it,dict(color=color))]
        para(tf,runs,size=size,line=line,sa=sa)

# ══════════════════════════════════════════════════════════════════════════
# 1 · COVER
s=slide()
rect(s,7.9,0.7,4.9,4.9,ln=GREEN,lw=1.5,shape=MSO_SHAPE.OVAL)
rect(s,9.2,3.4,3.2,3.2,ln=LINE,lw=1.0,shape=MSO_SHAPE.OVAL)
rect(s,7.62,1.02,0.32,0.32,fill=GREEN,shape=MSO_SHAPE.OVAL)
eyebrow(s,"Technical Overview · Engineering Audience")
tf=tf_box(s,MX,1.4,7.6,1.7)
para(tf,[("Application Architecture",dict(size=34,bold=True,color=BLACK))],line=1.05,sa=0)
para(tf,[("Three surfaces. One engine",dict(size=34,bold=True,color=BLACK)),(".",dict(size=34,bold=True,color=GREEN))],line=1.05,sa=0)
rect(s,MX,2.95,0.95,0.08,fill=GREEN)
tf=tf_box(s,MX,3.25,7.1,1.3)
para(tf,"H2A converts a SAP Hybris commerce monolith — Java/Spring, items.xml, and a Spartacus storefront — into a deployable Salesforce project of Apex, LWC, and metadata, with a human review gate at every consequential step.",size=13.5,color=GRAY,line=1.3)
fx=MX
for t,c,arrow in [("SAP HYBRIS",COPPER,True),("H2A ENGINE",BLACK,True),("SALESFORCE",TEAL,False)]:
    w=0.3+len(t)*0.095
    rect(s,fx,4.85,w,0.42,ln=c,lw=1.5)
    tf=tf_box(s,fx,4.85,w,0.42,anchor=MSO_ANCHOR.MIDDLE,wrap=False)
    para(tf,[(t,dict(size=10,bold=True,color=c,name=MONO))],align=PP_ALIGN.CENTER)
    fx+=w
    if arrow:
        tf=tf_box(s,fx,4.85,0.5,0.42,anchor=MSO_ANCHOR.MIDDLE)
        para(tf,[("──▶",dict(size=11,color=FAINT))],align=PP_ALIGN.CENTER); fx+=0.5
chips(s,[{'t':"Python 3.12 · FastAPI",'runs':[("Python 3.12 ",{}),("· FastAPI",dict(bold=True,color=BLACK))]},
         {'t':"React 18 + Vite · TypeScript",'runs':[("React 18 + Vite ",{}),("· TypeScript",dict(bold=True,color=BLACK))]},
         {'t':"85 automated tests",'runs':[("85",dict(bold=True,color=BLACK)),(" automated tests",{})]}],5.7)
footer(s,1)

# ══════════════════════════════════════════════════════════════════════════
# 2 · SYSTEM CONTEXT
s=slide()
header(s,"01 · System Context",["Everything is a client of one function"],2,tsize=28)
tf=tf_box(s,MX,BODY,CW,0.6)
para(tf,[("run_agentic_migration()",dict(name=MONO,bold=True,color=BLACK)),
         (" in orchestrator.py is the single entry point. The CLI, the VS Code extension, and the web "
          "cockpit are three ways to call it — they differ only in which of its optional hooks they pass.",{})],
     size=12.5,color=GRAY,line=1.3)

lx,cw=cols(3,0.3)
surf=[("CLI",TEAL,"python -m src.main agent-migrate",
       ["Direct, in-process call","Also: ingest · repo-migrate · impex · cronjob · metadata · report","The reference implementation"]),
      ("VS CODE EXTENSION",TEAL,"TypeScript + webview",
       ["Spawns the CLI as a child process","Bootstraps its own .venv","Ships a synced copy of the engine"]),
      ("WEB COCKPIT",TEAL,"React SPA + FastAPI",
       ["Imports the engine in-process","Runs it on a worker thread","Which is what makes live events and blocking gates possible"])]
for i,(name,col,sub,pts) in enumerate(surf):
    x=lx[i]; card(s,x,BODY+0.75,cw,2.45,top=col)
    tf=tf_box(s,x+0.24,BODY+1.0,cw-0.48,0.34)
    para(tf,[(name,dict(size=12.5,bold=True,color=BLACK,name=MONO,track=0.5))])
    tf=tf_box(s,x+0.24,BODY+1.36,cw-0.48,0.3)
    para(tf,[(sub,dict(size=9.5,color=col,name=MONO))])
    tf=tf_box(s,x+0.24,BODY+1.72,cw-0.48,1.5)
    bullets(tf,pts,dot=col,size=10.5,line=1.2,sa=5)

rect(s,MX,BODY+3.42,CW,0.5,fill=BLACK)
tf=tf_box(s,MX+0.2,BODY+3.42,CW-0.4,0.5,anchor=MSO_ANCHOR.MIDDLE)
para(tf,[("run_agentic_migration(input, output, *, on_event, gate, should_cancel, on_blackboard, state_dir)",
          dict(size=11,bold=True,color=WHITE,name=MONO))],align=PP_ALIGN.CENTER)

hooks=[("on_event","streams every stage, decision and artifact to the UI"),
       ("gate","suspends the run and waits for a human verdict"),
       ("should_cancel","cooperative cancellation at every safe point"),
       ("on_blackboard","exposes live state for diff, regenerate, Copilot")]
hx,hw=cols(4,0.22)
for i,(h,desc) in enumerate(hooks):
    rect(s,hx[i],BODY+4.12,hw,0.72,fill=PANEL,ln=LINE,lw=0.75)
    tf=tf_box(s,hx[i]+0.14,BODY+4.22,hw-0.28,0.24)
    para(tf,[(h,dict(size=10,bold=True,color=GDEEP,name=MONO))])
    tf=tf_box(s,hx[i]+0.14,BODY+4.46,hw-0.28,0.34)
    para(tf,[(desc,dict(size=8.5,color=GRAY))],line=1.14)

strip(s,"Why it matters",[("Every hook defaults to ",{}),("None",dict(name=MONO,bold=True,color=BLACK)),
    (" — so the engine has no knowledge of HTTP, of VS Code, or of a browser. That is what keeps all "
     "three surfaces at feature parity without a shared UI layer, and why adding the cockpit changed "
     "no CLI behaviour.",{})])

# ══════════════════════════════════════════════════════════════════════════
# 3 · THE DIAGRAM (full-bleed)
s=slide()
eyebrow(s,"02 · Architecture at a Glance")
tf=tf_box(s,MX,0.86,CW,0.5)
para(tf,[("Five layers, two rails",dict(size=26,bold=True,color=BLACK)),(".",dict(size=26,bold=True,color=GREEN))],line=1.0,sa=0)
png=DOCS/"architecture-diagram.png"
if png.exists():
    iw,ih=2600,1906
    h=6.15; w=h*(iw/ih)
    if w>CW: w=CW; h=w*(ih/iw)
    s.shapes.add_picture(str(png),Inches((PW-w)/2),Inches(1.48),Inches(w),Inches(h))
else:
    tf=tf_box(s,MX,3.5,CW,0.5)
    para(tf,[("architecture-diagram.png missing — run build_arch_image.py first",
              dict(size=13,color=RISK,name=MONO))],align=PP_ALIGN.CENTER)
footer(s,3)

# ══════════════════════════════════════════════════════════════════════════
# 4 · THE FIVE LAYERS
s=slide()
header(s,"03 · The Five Layers",["What each layer owns, and what it must not know"],4,tsize=26)
layers=[("SURFACES",TEAL,"CLI · VS Code extension · Web cockpit",
         "Know how to talk to a human. Know nothing about migration."),
        ("ORCHESTRATION",GDEEP,"orchestrator.py",
         "Stages, dependency wavefronts, human gates, checkpoints, failure containment."),
        ("AGENTS",GDEEP,"Blackboard · Planner · Builder · Critic · Verifier",
         "Agents never call each other — they read and write one shared workspace."),
        ("CAPABILITIES",BLACK,"ingest · schema · generate · validate · verify · rule_ledger · report",
         "Stateless functions. Independently testable. Shared by both pipelines."),
        ("LLM GATEWAY",PURPLE,"llm.py → anthropic · openrouter · mock",
         "Routing, retries, caching and cost accounting — one choke point, no bypass.")]
ly=BODY+0.12; lh=0.93
for i,(name,lc,nodes,why) in enumerate(layers):
    rect(s,MX,ly,CW,lh,fill=WHITE,ln=LINE,lw=0.75)
    rect(s,MX,ly,1.95,lh,fill=lc)
    tf=tf_box(s,MX+0.16,ly,1.7,lh,anchor=MSO_ANCHOR.MIDDLE)
    para(tf,[(name,dict(size=8.5,bold=True,color=WHITE,name=MONO,track=0.8))],line=1.1)
    tf=tf_box(s,MX+2.15,ly+0.17,CW-2.4,0.3)
    para(tf,[(nodes,dict(size=10.5,bold=True,color=BLACK,name=MONO))])
    tf=tf_box(s,MX+2.15,ly+0.5,CW-2.4,0.34)
    para(tf,[(why,dict(size=10,color=GRAY))],line=1.14)
    if i<len(layers)-1:
        tf=tf_box(s,MX,ly+lh-0.03,CW,0.22)
        para(tf,[("▼",dict(size=8,color=FAINT))],align=PP_ALIGN.CENTER)
    ly+=lh+0.19
footer(s,4)

# ══════════════════════════════════════════════════════════════════════════
# 5 · AGENTIC CORE
s=slide(PANEL)
header(s,"04 · The Agentic Core",["A shared blackboard, not a call graph"],5,tsize=27)
rect(s,MX,BODY+0.15,3.1,3.55,fill=BLACK)
tf=tf_box(s,MX+0.26,BODY+0.42,2.6,0.5)
para(tf,[("BLACKBOARD",dict(size=13,bold=True,color=WHITE,name=MONO,track=1.0))])
rect(s,MX+0.26,BODY+0.85,0.7,0.06,fill=GREEN)
tf=tf_box(s,MX+0.26,BODY+1.05,2.6,2.4)
for it in ["schema","comprehensions","plan (PlanItem)","artifacts (Artifact)","decisions","open questions"]:
    para(tf,[("· ",dict(color=GREEN)),(it,dict(color=RGBColor(0xc8,0xcb,0xce),name=MONO))],size=10.5,line=1.2,sa=7)

ax=MX+3.45; aw=(CW-3.45-0.3*2)/3
agents=[("planner.py",GDEEP,"Convert vs Skip per target. A native-product fit becomes a review flag — never a dropped class."),
        ("builders.py",GDEEP,"Builder + Verifier. Branches on layer: Java → Apex, Angular Component → LWC bundle."),
        ("critic.py",GDEEP,"Adversarial review of each artifact. Findings carry severity, category and a concrete fix."),
        ("retriever.py",PURPLE,"Lexical RAG over 8 bundled Salesforce docs. No vector DB — zero infra, zero drift."),
        ("router.py",PURPLE,"Cheap tier for comprehend and plan; frontier tier for generate, repair and critic."),
        ("incremental.py",BLACK,"Fingerprints source, deps, schema and plan so an unchanged class is never re-generated.")]
for i,(name,col,desc) in enumerate(agents):
    r,c=divmod(i,3)
    x=ax+c*(aw+0.3); y=BODY+0.15+r*1.85
    card(s,x,y,aw,1.62,top=col)
    tf=tf_box(s,x+0.2,y+0.24,aw-0.4,0.3)
    para(tf,[(name,dict(size=11.5,bold=True,color=BLACK,name=MONO))])
    tf=tf_box(s,x+0.2,y+0.6,aw-0.4,0.95)
    para(tf,[(desc,dict(size=10,color=GRAY))],line=1.22)
strip(s,"Why a blackboard",[("Because every intermediate state stays inspectable. That is precisely what makes the human "
    "review gates, the per-file diff, the regenerate-one-file action and the full audit trail possible at all — "
    "none of which you can bolt onto agents that call each other directly.",{})])

# ══════════════════════════════════════════════════════════════════════════
# 6 · THE PIPELINE
s=slide()
header(s,"05 · The Pipeline",["Six stages, three human gates"],6,tsize=28)
steps=[("1","Analyze",COPPER,"Ingest Java + Angular, parse items.xml, build the dependency graph and domain schedule.","Deterministic — no model calls"),
       ("2","Comprehend",BLACK,"One structured LLM pass per class: purpose, business rules, queries, side effects, risks.","Fully parallel"),
       ("3","Plan",COPPER,"Group classes into Salesforce targets; decide Convert vs Skip; attach review suggestions.","Cheap model tier"),
       ("4","Build + Critic",COPPER,"Generate Apex/LWC, validate, repair, then adversarially review each artifact.","One wavefront at a time"),
       ("5","Reconcile",BLACK,"Merge artifacts, augment schema, write the SFDX tree, emit data, schedules and reports.","Deterministic merge"),
       ("6","Verify",BLACK,"Optional real sf dry-run deploy; compiler errors feed the repair loop and redeploy.","Needs a live org")]
sx,sw=cols(6,0.16)
for i,(n,name,col,desc,note) in enumerate(steps):
    gated = col==COPPER
    card(s,sx[i],BODY+0.3,sw,3.35,top=col,fill=(RGBColor(0xfb,0xf7,0xf2) if gated else WHITE))
    tf=tf_box(s,sx[i]+0.16,BODY+0.5,sw-0.32,0.24)
    para(tf,[(f"STAGE {n}",dict(size=8,bold=True,color=FAINT,name=MONO,track=1.0))])
    tf=tf_box(s,sx[i]+0.16,BODY+0.78,sw-0.32,0.42)
    para(tf,[(name,dict(size=12.5,bold=True,color=BLACK))],line=1.05)
    tf=tf_box(s,sx[i]+0.16,BODY+1.3,sw-0.32,1.5)
    para(tf,[(desc,dict(size=9.5,color=GRAY))],line=1.24)
    tf=tf_box(s,sx[i]+0.16,BODY+2.88,sw-0.32,0.5)
    para(tf,[(note,dict(size=8.5,color=FAINT,italic=True))],line=1.14)
    if gated:
        rect(s,sx[i]+0.16,BODY+3.36,sw-0.32,0.02,fill=COPPER)
gy=BODY+3.82
for i,(lbl,idx) in enumerate([("⏸ DISCOVERY GATE",0),("⏸ PLAN GATE",2),("⏸ BUILD GATE",3)]):
    tf=tf_box(s,sx[idx],gy,sw*1.6,0.26)
    para(tf,[(lbl,dict(size=8.5,bold=True,color=COPPER,name=MONO,track=0.6))])
strip(s,"The gate that costs nothing",[("The Discovery gate fires ",{}),("before a single LLM call is made",dict(bold=True,color=BLACK)),
    (". A reviewer sees the complete repository understanding — file tree, architecture, every class and the "
     "data model — and can walk away at zero cost. Verified: 0 model requests before approval.",{})])

# ══════════════════════════════════════════════════════════════════════════
# 7 · SUPERVISED RUN
s=slide()
header(s,"06 · A Supervised Run",["How a review gate actually works"],7,tsize=28)
lx,cw=cols(2,0.4)
flow=[("1","Reviewer uploads a .zip and enables supervised mode"),
      ("2","POST /api/runs → FastAPI spawns a worker thread"),
      ("3","Client begins polling GET /api/runs/{id} every 1.2 s"),
      ("4","Analyze completes — still zero model calls"),
      ("5","gate(\"discovery\") blocks the worker thread on an Event"),
      ("6","Cockpit renders the scan; reviewer approves"),
      ("7","Event.set() releases the thread; Comprehend begins"),
      ("8","Plan and Build gates repeat the same exchange"),
      ("9","run_complete carries ledger, rule_ledger, cost, decisions")]
card(s,lx[0],BODY+0.2,cw,4.2,top=BLACK)
tf=tf_box(s,lx[0]+0.26,BODY+0.5,cw-0.52,0.3)
para(tf,[("REQUEST LIFECYCLE",dict(size=11,bold=True,color=BLACK,name=MONO,track=0.8))])
tf=tf_box(s,lx[0]+0.26,BODY+0.92,cw-0.52,3.2)
for n,txt in flow:
    para(tf,[(f"{n}  ",dict(color=GREEN,bold=True,name=MONO)),(txt,dict(color=GRAY))],size=10.5,line=1.2,sa=7)

card(s,lx[1],BODY+0.2,cw,2.0,top=TEAL)
tf=tf_box(s,lx[1]+0.26,BODY+0.5,cw-0.52,0.3)
para(tf,[("WHY POLLING, NOT SSE",dict(size=11,bold=True,color=BLACK,name=MONO,track=0.8))])
tf=tf_box(s,lx[1]+0.26,BODY+0.92,cw-0.52,1.2)
para(tf,[("An SSE endpoint still exists, but the shipped client does not use it. Corporate proxies were "
          "silently killing idle connections while a supervised run sat quiet at a gate. The client polls "
          "a replay-by-index endpoint instead — less elegant, dramatically more likely to work.",{})],
     size=10.5,color=GRAY,line=1.24)

card(s,lx[1],BODY+2.45,cw,1.95,top=GDEEP)
tf=tf_box(s,lx[1]+0.26,BODY+2.75,cw-0.52,0.3)
para(tf,[("WHAT THE REVIEWER CAN DO",dict(size=11,bold=True,color=BLACK,name=MONO,track=0.8))])
tf=tf_box(s,lx[1]+0.26,BODY+3.15,cw-0.52,1.2)
bullets(tf,["Approve, or reject with written feedback",
            "Inspect any file side-by-side against its Java origin",
            "Regenerate a single class without re-running the migration",
            "Ask the Copilot to rework a target"],dot=GDEEP,size=10.5,line=1.18,sa=5)
strip(s,"The hard part",[("The engine runs on a worker thread and a gate blocks that thread while HTTP stays responsive. "
    "An early version waited with no timeout and held the global run lock forever — which made a stopped "
    "migration un-restartable. Gates now honour cancellation, and a new run auto-cancels the previous one.",{})])

# ══════════════════════════════════════════════════════════════════════════
# 8 · DATA FLOW
s=slide(PANEL)
header(s,"07 · What Goes In, What Comes Out",["Code is half the deliverable"],8,tsize=26)
boxw=3.15; gap=(CW-3*boxw)/2
ins=[("Java / Spring","services, DAOs, facades, jobs"),("items.xml","the Hybris type system"),
     ("Spartacus","*.component.ts + templates"),("ImpEx · cronjobs","seed data and schedules")]
outs=[("force-app/classes","Apex + a test class for each"),("force-app/objects","SObjects and fields"),
      ("force-app/lwc","LWC bundles + @AuraEnabled Apex"),("data/","records converted from ImpEx")]
evid=[("MIGRATION_PLAN.md","every target, decision and rationale"),
      ("BUSINESS_RULES.md","every rule traced to code and test"),
      ("FEASIBILITY_REPORT.md","per-artifact confidence + completeness"),
      ("PARITY · DATA · CRON","coverage, records, schedules")]
for i,(hdr,col,items) in enumerate([("IN — SAP HYBRIS",COPPER,ins),("OUT — SALESFORCE DX",TEAL,outs),
                                    ("OUT — EVIDENCE",GREEN,evid)]):
    x=MX+i*(boxw+gap)
    card(s,x,BODY+0.25,boxw,3.9,top=col)
    tf=tf_box(s,x+0.24,BODY+0.55,boxw-0.48,0.3)
    para(tf,[(hdr,dict(size=11.5,bold=True,color=BLACK,name=MONO,track=0.6))])
    yy=BODY+1.0
    for name,sub in items:
        tf=tf_box(s,x+0.24,yy,boxw-0.48,0.28)
        para(tf,[(name,dict(size=10.5,bold=True,color=BLACK,name=MONO))])
        tf=tf_box(s,x+0.24,yy+0.26,boxw-0.48,0.3)
        para(tf,[(sub,dict(size=9.5,color=GRAY))],line=1.14)
        yy+=0.74
    if i<2:
        tf=tf_box(s,x+boxw,BODY+2.0,gap,0.4,anchor=MSO_ANCHOR.MIDDLE)
        para(tf,[("▶",dict(size=12,color=FAINT))],align=PP_ALIGN.CENTER)
strip(s,"The differentiator",[("Any tool can emit Apex. ",{}),("BUSINESS_RULES.md",dict(name=MONO,bold=True,color=BLACK)),
    (" measures completeness in business rules preserved rather than files converted — and names the rules "
     "that no generated artifact carries at all. That is the row nothing else reports.",{})])

# ══════════════════════════════════════════════════════════════════════════
# 9 · TECHNOLOGY STACK
s=slide()
header(s,"08 · Technology Stack",["What each piece is, and why it is there"],9,tsize=27)
groups=[("ENGINE",GDEEP,[("Python 3.12","host language for the pipeline"),
                         ("javalang","real Java AST — classes, methods, annotations"),
                         ("pydantic · pyyaml","structured-output validation; config + mappings"),
                         ("rich","CLI rendering")]),
        ("AI",PURPLE,[("anthropic SDK","primary provider; per-stage thinking effort"),
                      ("openai SDK → OpenRouter","alternate provider, OpenAI-compatible"),
                      ("bundled lexical RAG","8 Salesforce docs, top-3 per prompt — no vector DB"),
                      ("mock provider","keyless, deterministic; CI and safe demos")]),
        ("WEB",TEAL,[("FastAPI + uvicorn","15 JSON routes + SPA catch-all"),
                     ("React 18 + TypeScript 5.5","the cockpit, typed to the event contract"),
                     ("Vite 5 + Monaco","lazy-chunked diff editor; ~61 kB gzipped main bundle"),
                     ("markdown","renders .md reports to HTML in-app")]),
        ("DELIVERY",BLACK,[("VS Code API","webview + self-bootstrapping venv"),
                           ("Salesforce CLI (sf)","dry-run deploy — the only objective gate"),
                           ("Docker multi-stage → Render","node builds the SPA, python serves it"),
                           ("pytest","85 tests across pipeline, planner, LWC, rule ledger")])]
gx,gw=cols(2,0.34)
for i,(name,col,rows) in enumerate(groups):
    r,c=divmod(i,2)
    x=gx[c]; y=BODY+0.2+r*2.35
    card(s,x,y,gw,2.1,top=col)
    tf=tf_box(s,x+0.24,y+0.24,gw-0.48,0.3)
    para(tf,[(name,dict(size=11,bold=True,color=BLACK,name=MONO,track=0.9))])
    yy=y+0.62
    for tech,why in rows:
        tf=tf_box(s,x+0.24,yy,2.35,0.3)
        para(tf,[(tech,dict(size=9.5,bold=True,color=col,name=MONO))],line=1.1)
        tf=tf_box(s,x+2.66,yy,gw-2.9,0.34)
        para(tf,[(why,dict(size=9.5,color=GRAY))],line=1.14)
        yy+=0.35
footer(s,9)

# ══════════════════════════════════════════════════════════════════════════
# 10 · ARCHITECTURAL DECISIONS
s=slide()
header(s,"09 · Architectural Decisions",["The choices that shaped everything else"],10,tsize=26)
decs=[("Blackboard over a message bus",GDEEP,
       "Agents read and write one shared workspace; the orchestrator decides who runs when.",
       "Every intermediate state stays inspectable — which is what the review gates need."),
      ("Optional hooks, never subclassing",GDEEP,
       "Every integration point is a keyword argument defaulting to None.",
       "The cockpit was added without touching a single CLI code path."),
      ("Convert everything, flag natives",COPPER,
       "A better native home (CPQ, Flow, Approvals) produces a review flag on converted code.",
       "Silently dropping logic is how migrations lose requirements until go-live."),
      ("Deterministic parallelism",BLACK,
       "Concurrency bounded by dependency wavefronts; results merge in fixed order.",
       "A 7× speed-up is worthless if it makes output non-reproducible."),
      ("One LLM choke point",PURPLE,
       "Routing, retries, caching and cost accounting all live in llm.py.",
       "No call site can accidentally bypass the retry or the cost meter."),
      ("Polling over SSE",TEAL,
       "The client polls a replay-by-index endpoint instead of holding a stream open.",
       "Chosen after SSE plus heartbeats still failed on a real locked-down network.")]
dx,dw=cols(3,0.28)
for i,(name,col,what,why) in enumerate(decs):
    r,c=divmod(i,3)
    x=dx[c]; y=BODY+0.2+r*2.3
    card(s,x,y,dw,2.05,top=col)
    tf=tf_box(s,x+0.22,y+0.26,dw-0.44,0.56)
    para(tf,[(name,dict(size=12,bold=True,color=BLACK))],line=1.1)
    tf=tf_box(s,x+0.22,y+0.92,dw-0.44,0.6)
    para(tf,[(what,dict(size=9.5,color=GRAY))],line=1.2)
    rect(s,x+0.22,y+1.56,0.5,0.04,fill=col)
    tf=tf_box(s,x+0.22,y+1.66,dw-0.44,0.4)
    para(tf,[(why,dict(size=9,color=FAINT,italic=True))],line=1.18)
footer(s,10)

out=DOCS/"ARCHITECTURE_DECK.pptx"
prs.save(str(out))
print(f"wrote {out.relative_to(DOCS.parent)}  ({len(prs.slides._sldIdLst)} slides)")
