#!/usr/bin/env python
"""Build CLIENT_DECK.pptx — five pictorial slides for a first client conversation.

Deliberately not the 10-slide executive deck. This one is for the meeting where nobody has
seen the product yet: the job is to leave them with a *picture* of what happens, not a
specification. So every slide is a diagram with a caption, the text is short enough to read
from the back of a room, and the honest caveat gets its own place rather than being buried.

Shares the palette and helpers of build_deck.py so the three decks look like one family.

Regenerate:  python docs/build_client_deck.py
"""
import pathlib

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── palette (same family as the other two decks) ──
BLACK = RGBColor(0x0E, 0x0E, 0x0E); INK = RGBColor(0x22, 0x24, 0x27)
GRAY = RGBColor(0x53, 0x56, 0x5A);  FAINT = RGBColor(0x8F, 0x92, 0x96)
LINE = RGBColor(0xD2, 0xD2, 0xCF);  PANEL = RGBColor(0xF3, 0xF3, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); GREEN = RGBColor(0x86, 0xBC, 0x25)
GDEEP = RGBColor(0x04, 0x6A, 0x38); TEAL = RGBColor(0x00, 0x7C, 0xB0)
COPPER = RGBColor(0xA9, 0x62, 0x2F); RISK = RGBColor(0xB3, 0x28, 0x2D)
PALE = RGBColor(0xEC, 0xF3, 0xE4)
SANS, MONO = "Arial", "Courier New"

PW, PH = 13.333, 7.5          # 16:9 — projectors and Teams calls, not A4 handouts
MX = 0.62
CW = PW - 2 * MX
N = 5

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(PW), Inches(PH)


def slide(bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.shadow.inherit = False
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    return s


def rect(s, l, t, w, h, fill=None, ln=None, lw=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         adj=0.06):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = adj
        except (IndexError, AttributeError):
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if ln is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = ln; sp.line.width = Pt(lw)
    return sp


def tf_box(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _track(run, pts):
    run._r.get_or_add_rPr().set("spc", str(int(pts * 100)))


def para(tf, runs, align=PP_ALIGN.LEFT, size=12, color=INK, bold=False, name=SANS,
         sa=3, line=1.12):
    p0 = tf.paragraphs[0]
    p = p0 if not p0.runs else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(sa); p.line_spacing = line
    if isinstance(runs, str):
        runs = [(runs, {})]
    for txt, st in runs:
        r = p.add_run(); r.text = txt; f = r.font
        f.size = Pt(st.get("size", size)); f.bold = st.get("bold", bold)
        f.italic = st.get("italic", False); f.name = st.get("name", name)
        f.color.rgb = st.get("color", color)
        if st.get("track"):
            _track(r, st["track"])
    return p


def header(s, eyebrow, title, n):
    tf = tf_box(s, MX, 0.46, CW, 0.3)
    para(tf, [(eyebrow.upper(), dict(size=10, bold=True, color=GDEEP, name=MONO, track=2.0))])
    tf = tf_box(s, MX, 0.78, CW, 0.62)
    para(tf, [(title, dict(size=27, bold=True, color=BLACK))], line=1.02)
    rect(s, MX, 1.52, 0.9, 0.05, fill=GREEN, shape=MSO_SHAPE.RECTANGLE)
    foot(s, n)


def foot(s, n):
    tf = tf_box(s, MX, PH - 0.44, CW / 2, 0.24)
    para(tf, [("H2A  ·  HYBRIS TO SALESFORCE", dict(size=8, color=FAINT, name=MONO, track=1.4))])
    tf = tf_box(s, MX + CW / 2, PH - 0.44, CW / 2, 0.24)
    para(tf, [(f"{n:02d} / {N:02d}", dict(size=8, color=FAINT, name=MONO))],
         align=PP_ALIGN.RIGHT)


def caption(s, text, y=None, color=GRAY):
    """The one line under a diagram that says what it means."""
    tf = tf_box(s, MX, y if y is not None else PH - 1.32, CW, 0.7)
    para(tf, text, size=13, color=color, align=PP_ALIGN.CENTER, line=1.35)


def chip(s, l, t, w, h, label, sub="", accent=GREEN, fill=WHITE, big=13, small=9.5):
    """A labelled box — the deck's basic unit."""
    rect(s, l, t, w, h, fill=fill, ln=LINE, lw=1.0)
    rect(s, l, t, w, 0.075, fill=accent, shape=MSO_SHAPE.RECTANGLE)
    tf = tf_box(s, l + 0.14, t + 0.26, w - 0.28, h - 0.36,
                anchor=MSO_ANCHOR.TOP if sub else MSO_ANCHOR.MIDDLE)
    para(tf, [(label, dict(size=big, bold=True, color=BLACK))],
         align=PP_ALIGN.CENTER, line=1.06)
    if sub:
        para(tf, [(sub, dict(size=small, color=GRAY))], align=PP_ALIGN.CENTER,
             line=1.22, sa=0)


def arrow(s, l, t, w=0.44, h=0.34, color=FAINT, glyph="▶"):
    tf = tf_box(s, l, t, w, h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(glyph, dict(size=15, color=color))], align=PP_ALIGN.CENTER)


def icon(s, cx, cy, d, glyph, fg=WHITE, bg=GDEEP, size=20):
    rect(s, cx - d / 2, cy - d / 2, d, d, fill=bg, shape=MSO_SHAPE.OVAL)
    tf = tf_box(s, cx - d / 2, cy - d / 2, d, d, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(glyph, dict(size=size, color=fg))], align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════ 1 · COVER ══
s = slide(BLACK)
rect(s, 0, 0, PW, 0.16, fill=GREEN, shape=MSO_SHAPE.RECTANGLE)

tf = tf_box(s, MX, 1.55, CW, 0.4)
para(tf, [("AI-POWERED MIGRATION ACCELERATOR", dict(size=11, bold=True, color=GREEN,
                                                    name=MONO, track=2.4))],
     align=PP_ALIGN.CENTER)

tf = tf_box(s, MX, 2.05, CW, 1.5)
para(tf, [("SAP Hybris", dict(size=42, bold=True, color=WHITE)),
          ("  →  ", dict(size=42, color=GREEN)),
          ("Salesforce", dict(size=42, bold=True, color=WHITE))],
     align=PP_ALIGN.CENTER, line=1.06)
para(tf, [("Migrated by an AI agent team. ", dict(size=17, color=RGBColor(0xC9, 0xCC, 0xD0))),
          ("Proven to still behave the same.", dict(size=17, bold=True, color=GREEN))],
     align=PP_ALIGN.CENTER, line=1.3, sa=10)

# the three-step picture, in the largest form it appears anywhere in the deck
y, bw, gap = 4.05, 3.15, 0.62
x = (PW - (3 * bw + 2 * gap)) / 2
for i, (label, sub, col, gl) in enumerate([
        ("Your Hybris code", "Java · data model · data · jobs · storefront", COPPER, "❯"),
        ("AI agent team", "plan · build · review · verify", GREEN, "◆"),
        ("Deployable Salesforce", "Apex · LWC · objects · Flows + the evidence", TEAL, "✓")]):
    rect(s, x, y, bw, 1.5, fill=RGBColor(0x1A, 0x1C, 0x1F), ln=RGBColor(0x33, 0x36, 0x3A), lw=1.0)
    rect(s, x, y, bw, 0.07, fill=col, shape=MSO_SHAPE.RECTANGLE)
    icon(s, x + bw / 2, y + 0.5, 0.42, gl, fg=BLACK, bg=col, size=15)
    tf = tf_box(s, x + 0.16, y + 0.78, bw - 0.32, 0.62)
    para(tf, [(label, dict(size=14, bold=True, color=WHITE))], align=PP_ALIGN.CENTER)
    para(tf, [(sub, dict(size=9.5, color=FAINT))], align=PP_ALIGN.CENTER, line=1.24, sa=0)
    if i < 2:
        arrow(s, x + bw + 0.09, y + 0.58, gap - 0.18, 0.34, color=GREEN, glyph="▶")
    x += bw + gap

tf = tf_box(s, MX, 6.15, CW, 0.4)
para(tf, [("Web cockpit  ·  VS Code  ·  Command line", dict(size=11, color=FAINT, name=MONO)),
          ("        Powered by Anthropic Claude", dict(size=11, color=FAINT, name=MONO))],
     align=PP_ALIGN.CENTER)
foot(s, 1)


# ═══════════════════════════════════════════════════════════ 2 · WHAT IT DOES ══
s = slide()
header(s, "01 · What it does", "Point it at the old system. Get a working new one — with receipts.", 2)

COLW = (CW - 1.5) / 3
xs = [MX, MX + COLW + 0.75, MX + 2 * (COLW + 0.75)]
top = 1.95

# IN
rect(s, xs[0], top, COLW, 3.5, fill=PANEL, ln=LINE, lw=1.0)
rect(s, xs[0], top, COLW, 0.075, fill=COPPER, shape=MSO_SHAPE.RECTANGLE)
tf = tf_box(s, xs[0] + 0.2, top + 0.3, COLW - 0.4, 0.4)
para(tf, [("IN — what you have", dict(size=13, bold=True, color=BLACK))])
tf = tf_box(s, xs[0] + 0.2, top + 0.82, COLW - 0.4, 2.5)
for it in ["Java business logic — orders, pricing, customers",
           "The data model (items.xml)",
           "Real data (ImpEx files)",
           "Scheduled jobs (cronjobs)",
           "Workflows (business processes)",
           "The Spartacus storefront (Angular)"]:
    para(tf, [("—  ", dict(color=COPPER)), (it, dict(color=GRAY))], size=11.5, line=1.25, sa=9)

# ENGINE
rect(s, xs[1], top, COLW, 3.5, fill=BLACK)
rect(s, xs[1], top, COLW, 0.075, fill=GREEN, shape=MSO_SHAPE.RECTANGLE)
icon(s, xs[1] + COLW / 2, top + 0.85, 0.72, "◆", fg=BLACK, bg=GREEN, size=26)
tf = tf_box(s, xs[1] + 0.2, top + 1.36, COLW - 0.4, 0.5)
para(tf, [("Four AI agents", dict(size=14, bold=True, color=WHITE))], align=PP_ALIGN.CENTER)
tf = tf_box(s, xs[1] + 0.24, top + 1.85, COLW - 0.48, 1.5)
for role, what in [("Planner", "decides what becomes what"),
                   ("Builder", "writes the code + tests"),
                   ("Critic", "reviews it, adversarially"),
                   ("Verifier", "deploys it for real")]:
    para(tf, [(f"{role}  ", dict(size=11.5, bold=True, color=GREEN)),
              (what, dict(size=11, color=RGBColor(0xB8, 0xBC, 0xC0)))], line=1.2, sa=7)

# OUT
rect(s, xs[2], top, COLW, 3.5, fill=PALE, ln=LINE, lw=1.0)
rect(s, xs[2], top, COLW, 0.075, fill=TEAL, shape=MSO_SHAPE.RECTANGLE)
tf = tf_box(s, xs[2] + 0.2, top + 0.3, COLW - 0.4, 0.4)
para(tf, [("OUT — what you get", dict(size=13, bold=True, color=BLACK))])
tf = tf_box(s, xs[2] + 0.2, top + 0.82, COLW - 0.4, 2.5)
for it in ["Apex classes + a test for every one",
           "Custom objects and fields",
           "Lightning Web Components",
           "Salesforce Flows from your workflows",
           "Data loaders and job schedules",
           "**The evidence it still works**"]:
    bold = it.startswith("**")
    txt = it.strip("*")
    para(tf, [("—  ", dict(color=TEAL)),
              (txt, dict(color=BLACK if bold else GRAY, bold=bold))],
         size=11.5, line=1.25, sa=9)

arrow(s, xs[0] + COLW + 0.16, top + 1.55, 0.43, 0.4, color=FAINT)
arrow(s, xs[1] + COLW + 0.16, top + 1.55, 0.43, 0.4, color=FAINT)

caption(s, [("Months of manual rewriting becomes a reviewed first draft in ", {}),
            ("minutes to hours", dict(bold=True, color=BLACK)),
            (" — and every business rule is tracked from the old code to the new.", {})],
        y=5.72)


# ══════════════════════════════════════════════════════════ 3 · ARCHITECTURE ══
s = slide()
header(s, "02 · Architecture", "Three front doors. One engine. One place the AI is called.", 3)

y = 1.92
LH, GAP = 0.86, 0.3

# L1 surfaces
rect(s, MX, y, CW, LH, fill=WHITE, ln=LINE, lw=1.0)
tf = tf_box(s, MX + 0.22, y, 2.1, LH, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("YOU USE IT", dict(size=9.5, bold=True, color=TEAL, name=MONO, track=1.0))])
bx, bw2 = MX + 2.35, (CW - 2.6) / 3 - 0.14
for lbl, sub in [("Web cockpit", "watch and approve live"),
                 ("VS Code", "right-click a folder"),
                 ("Command line", "for CI/CD")]:
    chip(s, bx, y + 0.13, bw2, LH - 0.26, lbl, sub, accent=TEAL, big=11.5, small=8.5)
    bx += bw2 + 0.21
y += LH + GAP
arrow(s, PW / 2 - 0.2, y - GAP + 0.02, 0.4, GAP - 0.04, color=FAINT, glyph="▼")

# L2 engine
rect(s, MX, y, CW, LH * 1.28, fill=PANEL, ln=LINE, lw=1.0)
tf = tf_box(s, MX + 0.22, y, 2.1, LH * 1.28, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("THE ENGINE", dict(size=9.5, bold=True, color=GDEEP, name=MONO, track=1.0))])
bx, bw3 = MX + 2.35, (CW - 2.6) / 4 - 0.13
for lbl, sub, ac in [("Planner", "what becomes what", GREEN),
                     ("Builder", "writes Apex + LWC", GREEN),
                     ("Critic", "challenges it", GREEN),
                     ("Verifier", "deploys to your org", GREEN)]:
    chip(s, bx, y + 0.16, bw3, LH * 1.28 - 0.32, lbl, sub, accent=ac, big=11.5, small=8.5)
    bx += bw3 + 0.17
y += LH * 1.28 + GAP
arrow(s, PW / 2 - 0.2, y - GAP + 0.02, 0.4, GAP - 0.04, color=FAINT, glyph="▼")

# L3 assurance
rect(s, MX, y, CW, LH, fill=PALE, ln=GREEN, lw=1.0)
tf = tf_box(s, MX + 0.22, y, 2.1, LH, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("THE PROOF", dict(size=9.5, bold=True, color=GDEEP, name=MONO, track=1.0))])
tf = tf_box(s, MX + 2.35, y, CW - 2.6, LH, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("Rule ledger  ·  your own tests replayed  ·  line-level traceability  ·  "
           "review triage  ·  sign-off", dict(size=11.5, bold=True, color=GDEEP))])
para(tf, [("No AI is used here at all — it checks the AI's work by reading what the run "
           "recorded.", dict(size=9.5, color=GRAY))], sa=0)
y += LH + GAP
arrow(s, PW / 2 - 0.2, y - GAP + 0.02, 0.4, GAP - 0.04, color=FAINT, glyph="▼")

# L4 model
rect(s, MX, y, CW, LH * 0.8, fill=WHITE, ln=LINE, lw=1.0)
tf = tf_box(s, MX + 0.22, y, 2.1, LH * 0.8, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("THE AI", dict(size=9.5, bold=True, color=COPPER, name=MONO, track=1.0))])
tf = tf_box(s, MX + 2.35, y, CW - 2.6, LH * 0.8, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("One gateway  —  Anthropic Claude  ·  a cheaper model for simple work  ·  "
           "an offline mode that makes no calls at all",
           dict(size=11.5, color=INK))])

caption(s, [("Every AI call goes through one door, so cost, retries and caching are "
             "controlled in one place — and ", {}),
            ("the offline mode lets you evaluate the whole thing with zero exposure",
             dict(bold=True, color=BLACK)), (".", {})], y=6.32)


# ═════════════════════════════════════════════════════════════ 4 · WORKFLOW ══
s = slide()
header(s, "03 · How a migration runs", "Minutes of machine work. Three moments that are yours.", 4)

# the rail
railY = 3.16
rect(s, MX + 0.3, railY, CW - 0.6, 0.05, fill=LINE, shape=MSO_SHAPE.RECTANGLE)

STEPS = [
    ("Read", "Finds every file and\nworks out the shape", GDEEP, False),
    ("Check the bill", "What it will cost, and\nwhat clashes in your org", GDEEP, False),
    ("YOU APPROVE", "Before a single\nAI call is made", COPPER, True),
    ("Understand", "What each class does,\nand the rules inside it", GDEEP, False),
    ("YOU APPROVE", "The plan — change\nanything you disagree with", COPPER, True),
    ("Build & review", "Writes it, then a second\nAI challenges it", GDEEP, False),
    ("YOU APPROVE", "The code — or send it\nback for a rework", COPPER, True),
    ("Verify & prove", "Deploys for real, then\nproves the rules survived", GDEEP, False),
]
bw4 = (CW - 0.6) / len(STEPS)
for i, (title, sub, col, gate) in enumerate(STEPS):
    cx = MX + 0.3 + bw4 * i + bw4 / 2
    icon(s, cx, railY + 0.025, 0.38 if gate else 0.3, "▲" if gate else "●",
         fg=WHITE, bg=col, size=12 if gate else 9)
    if gate:
        # gates hang below the rail so the human moments read as a separate rhythm
        tf = tf_box(s, cx - bw4 / 2, railY + 0.4, bw4, 0.34)
        para(tf, [(title, dict(size=9.5, bold=True, color=COPPER, name=MONO, track=0.6))],
             align=PP_ALIGN.CENTER)
        tf = tf_box(s, cx - bw4 / 2 + 0.04, railY + 0.76, bw4 - 0.08, 0.9)
        para(tf, [(sub, dict(size=9, color=GRAY))], align=PP_ALIGN.CENTER, line=1.3)
    else:
        tf = tf_box(s, cx - bw4 / 2, railY - 1.04, bw4, 0.34)
        para(tf, [(title, dict(size=12, bold=True, color=BLACK))], align=PP_ALIGN.CENTER)
        tf = tf_box(s, cx - bw4 / 2 + 0.04, railY - 0.66, bw4 - 0.08, 0.62)
        para(tf, [(sub, dict(size=9, color=GRAY))], align=PP_ALIGN.CENTER, line=1.28)

# the two bands that explain the rhythm
rect(s, MX, 1.72, CW, 0.3, fill=PALE, shape=MSO_SHAPE.RECTANGLE)
tf = tf_box(s, MX + 0.16, 1.72, CW - 0.32, 0.3, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("THE MACHINE WORKS", dict(size=9, bold=True, color=GDEEP, name=MONO, track=1.2))])

rect(s, MX, 4.86, CW, 0.32, fill=RGBColor(0xF7, 0xEF, 0xE7), shape=MSO_SHAPE.RECTANGLE)
tf = tf_box(s, MX + 0.16, 4.86, CW - 0.32, 0.32, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("YOU DECIDE  —  AND CAN STOP AT ANY POINT TO FIX THE SOURCE, THEN RE-RUN",
           dict(size=9, bold=True, color=COPPER, name=MONO, track=1.2))])

caption(s, [("Nothing happens that you did not approve. ", dict(bold=True, color=BLACK)),
            ("The first gate comes ", {}),
            ("before any AI is used", dict(bold=True, color=BLACK)),
            (" — so if the tool has misread your codebase, you find out at zero cost.", {})],
        y=5.5)


# ══════════════════════════════════════════════════════ 5 · WHY IT'S TRUSTED ══
s = slide(PANEL)
header(s, "04 · Why you can trust the output", "Anyone can hand you Apex. This proves it still behaves the same.", 5)

y, cw5 = 1.95, (CW - 0.66) / 3
for i, (glyph, title, body) in enumerate([
        ("◎", "Every rule is tracked",
         "Each business rule found in your Java is followed to the finished code and "
         "lands in one of four buckets — including “nothing carries this any more”. "
         "That last one is what no other tool shows you."),
        ("↺", "Your own tests, replayed",
         "We take your existing JUnit suite, extract what the old code actually did, "
         "and replay it against the new Apex. The AI can set a test up — it is never "
         "allowed to write the expected answer."),
        ("✓", "Deployed, not assumed",
         "Connect an org and the output is deployed to it — validation only, nothing "
         "destructive. Real compiler errors are read and fixed until it is green, and "
         "the sign-off says plainly whether this step ran.")]):
    x = MX + i * (cw5 + 0.33)
    rect(s, x, y, cw5, 2.72, fill=WHITE, ln=LINE, lw=1.0)
    rect(s, x, y, cw5, 0.08, fill=GREEN, shape=MSO_SHAPE.RECTANGLE)
    icon(s, x + cw5 / 2, y + 0.72, 0.6, glyph, fg=WHITE, bg=GDEEP, size=21)
    tf = tf_box(s, x + 0.22, y + 1.16, cw5 - 0.44, 0.42)
    para(tf, [(title, dict(size=13.5, bold=True, color=BLACK))], align=PP_ALIGN.CENTER)
    tf = tf_box(s, x + 0.22, y + 1.62, cw5 - 0.44, 1.0)
    para(tf, [(body, dict(size=10.5, color=GRAY))], align=PP_ALIGN.CENTER, line=1.34)

# the honesty band — given its own weight on purpose
rect(s, MX, 5.02, CW, 1.12, fill=WHITE, ln=COPPER, lw=1.25)
rect(s, MX, 5.02, 0.075, 1.12, fill=COPPER, shape=MSO_SHAPE.RECTANGLE)
tf = tf_box(s, MX + 0.34, 5.2, CW - 0.6, 0.34)
para(tf, [("AND IT TELLS YOU WHAT IT CANNOT PROVE",
           dict(size=10, bold=True, color=COPPER, name=MONO, track=1.6))])
tf = tf_box(s, MX + 0.34, 5.55, CW - 0.6, 0.5)
para(tf, [("Every run ends with a sign-off document recording who approved what, on what "
           "evidence — and, listed just as prominently, the rules with no test, the code "
           "with no traceable origin, and whether a real org ever compiled it. ", {}),
          ("There is no “100%” badge anywhere in this product, deliberately.",
           dict(bold=True, color=BLACK))], size=11.5, line=1.32)

caption(s, [("A number you can trust when it is green ", dict(bold=True, color=BLACK)),
            ("is worth more than a green number.", {})], y=6.42, color=GRAY)


out = pathlib.Path(__file__).resolve().parent / "CLIENT_DECK.pptx"
prs.save(str(out))
print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
